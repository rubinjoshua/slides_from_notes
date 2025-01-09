import os
from parser import parse_file
from pptx.util import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import random
import math
from svgpathtools import svg2paths
import shutil
import re
from tqdm import tqdm
import hashlib
from image_generator import generate_stable_diffusion_image
from collections import defaultdict

MD_FILE_PATH = "/Users/joshuarubin/Library/Mobile Documents/com~apple~CloudDocs/לימוד הלכה/שבת"
RABBI_FILE_PATH = "./Rabbi Stickers"
ALTERNATIVE_RABBI_PATH = "./Rabbi Stickers backup"
BUBBLE_SVG_PATH = "./up_bubble.svg"
IMAGES_PATH = "./images/"
CONTINUE_DATA = "./continue.txt"
NUM_OF_SLIDES = 1

def mkdir_if_not_exist(path):
    if not os.path.exists(path):
        os.mkdir(path)


class SimanPowerpoint():
    slide_width = Pt(1740)  # = width of 300ppi A6
    slide_height = Pt(1230)  # = height of 300ppi A6
    slide_bg_colour = RGBColor(255, 255, 255)

    first_rabbi_x = Pt(1462)
    first_rabbi_y = Pt(250)

    rabbi_height = Pt(220)
    small_rabbi_height = Pt(135)

    rabbis_together_offset = Pt(100)
    small_rabbis_together_offset = Pt(65)
    rabbis_apart_offset = Pt(325)
    rabbis_row_offset = Pt(450)

    chars_in_bubble = 40000

    huge_text_size = Pt(524)
    large_text_size = Pt(164)

    base_font = 'Secular One'

    image_margin = Pt(60)

    image_width = slide_width - image_margin
    image_height = slide_height - image_margin

    long_slide = 12

    def __init__(self):
        self.image_path = None
        self.prs = Presentation()
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height
        self.siman = ""
        self.current_slide = None
        self.add_slide()
        self.bubble_points = []
        for ij in svg2paths(BUBBLE_SVG_PATH)[0][0]:
            self.bubble_points.append((Pt(ij.start.real), Pt(ij.start.imag)))
            self.bubble_points.append((Pt(ij.end.real), Pt(ij.end.imag)))
        top_left_idx = self.bubble_points.index(min(self.bubble_points, key=lambda p: p[0] + p[1]))
        self.bubble_points = self.bubble_points[top_left_idx:] + self.bubble_points[:top_left_idx]
        self.rabbi_paths = {re.sub("[^א-ת ״]", "", f).strip(): os.path.join(RABBI_FILE_PATH, f)
                                        for f in os.listdir(RABBI_FILE_PATH) if not f.startswith(".") and "png" in f}
        self.alternative_rabbi_paths = {re.sub("[^א-ת ״]", "", f).strip(): os.path.join(ALTERNATIVE_RABBI_PATH, f)
                                        for f in os.listdir(ALTERNATIVE_RABBI_PATH) if not f.startswith(".") and "png" in f}
        self.undrawn_rabbis = []

    def __del__(self):
        self.prs.save(os.path.join(MD_FILE_PATH, self.siman + '.pptx'))
        if self.undrawn_rabbis:
            print("Un-drawn Rabbis:")
            for r in set(self.undrawn_rabbis):
                print(r)

    def add_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide_fill = slide.background.fill
        slide_fill.solid()
        slide_fill.fore_color.rgb = self.slide_bg_colour
        self.current_slide = slide

    def add_large_centered_text(self, text, is_seif):
        title_box = self.current_slide.shapes.add_textbox(0, 0, height=self.slide_height, width=self.slide_width)
        title_box.line.width = 0
        text_frame = title_box.text_frame
        text_frame.clear()

        # text_frame.text = text
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        font = run.font
        font.name = 'Noto Rashi Hebrew' if is_seif else self.base_font

        run.text = text
        # font = text_frame.paragraphs[0].font
        for paragraph in text_frame.paragraphs:
            for r in paragraph.runs:
                r.font.name = 'Noto Rashi Hebrew' if is_seif else self.base_font
        # for char in run.text:
        #     char_run = p.add_run()
        #     char_run.text = char
        #     char_run.font.name = 'Noto Rashi Hebrew' if is_seif else self.base_font

        font.bold = True
        font.size = self.huge_text_size if is_seif else self.large_text_size
        font.color.rgb = RGBColor(0, 0, 0)
        title_box.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        title_box.left = int((self.slide_width/2) - (title_box.width/2))
        title_box.top = int((self.slide_width/2) - (title_box.height/2))

    def add_corner_seif(self, seif):
        text = ", ".join([self.siman, seif])
        text_box = self.current_slide.shapes.add_textbox(Pt(18), Pt(1142), Pt(100), Pt(100))
        text_box.text_frame.text = text
        font = text_box.text_frame.paragraphs[0].font
        font.bold = True
        font.size = Pt(60)
        font.color.rgb = RGBColor(0, 0, 0)
        font.name = self.base_font
        text_box.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def add_question_slide(self, question, seif):
        self.add_slide()
        self.add_large_centered_text(question, is_seif=False)
        self.add_corner_seif(seif)

    def add_rabbi(self, rabbi_path, x, y, h):
        rabbi = self.current_slide.shapes.add_picture(rabbi_path, x, y, height=h)
        rabbi.rotation = (random.random() * 4) - 2
        return int(rabbi.width.pt)

    def add_rabbi_name(self, x, y, h, w, name):
        self.undrawn_rabbis.append(name)
        name_box = self.current_slide.shapes.add_textbox(x, y, height=h, width=w)
        name_box.line.width = 0
        text_frame = name_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = name
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        font = text_frame.paragraphs[0].font
        font.name = 'Helvetica'
        font.size = Pt(24)
        font.color.rgb = RGBColor(0, 0, 0)
        name_box.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        return int(name_box.width.pt)

    def add_rabbi_or_name(self, file_name, rx, ry, small=False):
        file_path = self.rabbi_paths.get(file_name)
        alternative_file_path = self.alternative_rabbi_paths.get(file_name)
        if small:
            rx -= Pt(15)
            ry -= Pt(15)
        height = self.small_rabbi_height if small else self.rabbi_height
        if alternative_file_path and not file_path:
            file_path = os.path.join(RABBI_FILE_PATH, alternative_file_path.split("/")[-1])
            shutil.copyfile(alternative_file_path, file_path)
        if file_path:
            return self.add_rabbi(file_path, rx, ry, height)
        return self.add_rabbi_name(rx, ry + (Pt(120) if not small else Pt(40)), Pt(20), Pt(150), file_name)

    def update_rx(self, rx, is_last, is_small, is_together):
        if is_last:
            return rx
        if is_together:
            offset = self.small_rabbis_together_offset if is_small else self.rabbis_together_offset
        else:
            offset = self.rabbis_apart_offset
        return rx - offset

    def add_group_of_rabbis(self, rx, ry, rabbis):
        right_side_of_group_of_rabbis = rx
        for i, rabbi_refs in enumerate(rabbis):
            rabbi_width = self.add_rabbi_or_name(rabbi_refs["rabbi"], rx, ry, small=False)
            if i == 0:
                right_side_of_group_of_rabbis += Pt(rabbi_width)
            for j, ref in enumerate(rabbi_refs["refs"]):
                self.add_rabbi_or_name(ref, rx, ry, small=True)
                rx = self.update_rx(rx, is_last=False, is_small=True, is_together=True)
            rx = self.update_rx(rx, is_last=i == len(rabbis) - 1, is_small=False, is_together=True)
        rabbi_group_width = right_side_of_group_of_rabbis - rx
        return rabbi_group_width, rx

    def add_images(self, images_folder):
        for image in [i for i in os.listdir(self.image_path) if i.startswith(images_folder) and i.endswith("_t.png")]:
            i_x = int(self.slide_width/2) - int(self.image_width/2)
            i_y = int(self.slide_height/2) - int(self.image_height/2)
            image = self.current_slide.shapes.add_picture(os.path.join(self.image_path, image), i_x, i_y,
                                                          height=self.image_height, width=self.image_width)
            image.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

    def add_bubble(self, text, bx, by, rabbi_group_width):
        scale = math.sqrt(len(text) / self.chars_in_bubble)
        points = [(p[0] * scale, p[1] * scale) for p in self.bubble_points]
        bubble_builder = self.current_slide.shapes.build_freeform(start_x=points[0][0], start_y=points[0][1])
        bubble_builder.add_line_segments(points)
        bubble = bubble_builder.convert_to_shape()
        bubble.left = bx - int(bubble.width / 2) + int(rabbi_group_width / 2)
        bubble.top = by - bubble.height + Pt(52)
        bubble.shadow.inherit = False

        fill = bubble.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        fill.transparency = 0.12

        line = bubble.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(6)

        text_frame = bubble.text_frame
        edge_margin = Pt(100 * scale)

        text_frame.margin_top = edge_margin
        text_frame.margin_bottom = Pt(800 * scale)
        text_frame.margin_right = edge_margin
        text_frame.margin_left = edge_margin

        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        font = p.font
        font.name = 'Helvetica'
        font.size = Pt(24)
        font.color.rgb = RGBColor(0, 0, 0)

    def add_responses_slide(self, responses, images_folder):
        self.add_slide()
        self.add_images(images_folder)
        rx = self.first_rabbi_x
        ry = self.first_rabbi_y
        for response_idx, response in responses.items():
            if response_idx != 0 and response_idx % 5 == 0:
                ry += self.rabbis_row_offset
                rx = self.first_rabbi_x
            rabbi_group_width, rx = self.add_group_of_rabbis(rx, ry, response["rabbis"])
            if "text" in response.keys():
                self.add_bubble(response["text"], rx, ry, rabbi_group_width)
            rx = self.update_rx(rx, is_last=False, is_small=False, is_together=False)

    def get_images_folder(self, question):
        return hashlib.md5(question.encode()).hexdigest()

    def parse_data_to_slides(self, slides_data):
        self.siman = slides_data["siman"]
        self.image_path = os.path.join(IMAGES_PATH, self.siman)
        mkdir_if_not_exist(self.image_path)
        self.add_large_centered_text(self.siman, is_seif=True)
        self.add_slide() # this is just a blank slide
        pbar = tqdm(total=sum([1 for v in slides_data.values() if type(v) == dict for v2 in v.values()]), desc="Generating")
        for seif, seif_data in [(k, v) for k, v in slides_data.items() if type(v) == dict]:
            for question, responses in seif_data.items():
                pbar.update()
                self.add_question_slide(question, seif)
                images_folder = self.get_images_folder(question)
                caption = responses.pop("caption")
                generate_stable_diffusion_image(
                    path=os.path.join(self.image_path, images_folder),
                    text=caption,
                    fake_gen=not actually_generate_image)
                if len(responses) > self.long_slide:
                    fst_half_of_rabbis = {k: v for k, v in responses.items() if k < round(len(responses)/2)}
                    snd_half_of_rabbis = {k: v for k, v in responses.items() if k >= round(len(responses)/2)}
                    fst_half_of_rabbis[max(fst_half_of_rabbis.keys()) + 1] = {"rabbis": [{"rabbi": "המשך א", "refs": []}]}
                    snd_half_of_rabbis[max(snd_half_of_rabbis.keys()) + 1] = {"rabbis": [{"rabbi": "המשך ב", "refs": []}]}
                    snd_half_of_rabbis = {i: v for i, (k, v) in enumerate(snd_half_of_rabbis.items())}
                    self.add_responses_slide(fst_half_of_rabbis, images_folder)
                    self.add_question_slide("(המשך)", seif)
                    self.add_responses_slide(snd_half_of_rabbis, images_folder)
                else:
                    self.add_responses_slide(responses, images_folder)
        pbar.close()


actually_generate_image = True  # each slide costs 3.93c to generate (3.93 DreamStudio credits)


def get_latest_location():
    if os.path.exists(CONTINUE_DATA):
        s, n = open(CONTINUE_DATA, "r+").readlines()
        n = int(n)
        s = s.strip()
    else:
        s = file_names_idxs[min(file_names_idxs.keys())]
        n = 0
    return s, n


def get_last_n_from_slides_data(slides_data, s):
    small_slides_data = {"siman": slides_data.pop("siman")}

    seif_questions = [(seif, (question, seif_data[question])) for seif, seif_data in slides_data.items() for question in
                      seif_data]

    end_n = n + NUM_OF_SLIDES
    if end_n > len(seif_questions):
        end_n = len(seif_questions)
        s_k = [k for k, v in file_names_idxs.items() if v == s][0]
        s = file_names_idxs[str(int(s_k) + 1)]
        next_n = "0"
    else:
        next_n = str(end_n)

    seif_questions = seif_questions[n:end_n]
    n_data = defaultdict(dict)
    for seif_q in seif_questions:
        n_data[seif_q[0]][seif_q[1][0]] = seif_q[1][1]
    small_slides_data = small_slides_data | n_data
    slides_data = small_slides_data
    return slides_data, s, next_n


def update_location_for_next_continue(s, n, next_n, satisfied):
    if satisfied:
        print("All rabbis present.")
        open(CONTINUE_DATA, "w+").write("\n".join([s, next_n]))
    else:
        open(CONTINUE_DATA, "w+").write("\n".join([s, str(n)]))


if __name__ == "__main__":
    file_names = sorted([os.path.join(MD_FILE_PATH, f) for f in os.listdir(MD_FILE_PATH) if f.endswith(".md")])
    file_names_idxs = {str(i+1): f for i, f in enumerate(file_names)}

    s, n = get_latest_location()
    file_name = s
    # print("Files:")
    # for i, f in file_names_idxs.items():
    #     print(i, "-", f.split("/")[-1].split(".md")[0])
    # file_name = file_names_idxs[str(input("Which file?"))]

    presentation = SimanPowerpoint()
    slides_data = parse_file(file_name)

    slides_data, s, next_n = get_last_n_from_slides_data(slides_data, s)

    presentation.parse_data_to_slides(slides_data)

    update_location_for_next_continue(s, n, next_n, satisfied=not presentation.undrawn_rabbis)

    del presentation




